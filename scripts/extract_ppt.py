#!/usr/bin/env python3
"""
Extract text content from PowerPoint (.pptx) files.
Handles Chinese character encoding properly.

Usage:
    python extract_ppt.py <file.pptx> [--output json|text] [--slide N]
"""

import sys
import json
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def extract_slide_text(slide):
    """Extract all text from a single slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
    return texts


def extract_slide_tables(slide):
    """Extract tables from a slide."""
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            tables.append(table_data)
    return tables


def extract_pptx(file_path, slide_num=None):
    """
    Extract content from a PPTX file.

    Args:
        file_path: Path to the .pptx file
        slide_num: Optional specific slide number to extract (1-based)

    Returns:
        List of dictionaries with slide content
    """
    prs = Presentation(file_path)
    slides_content = []

    for i, slide in enumerate(prs.slides):
        slide_number = i + 1

        # Skip if specific slide requested and this isn't it
        if slide_num and slide_number != slide_num:
            continue

        slide_text = extract_slide_text(slide)
        slide_tables = extract_slide_tables(slide)

        slides_content.append({
            'slide_num': slide_number,
            'text': slide_text,
            'tables': slide_tables,
            'text_combined': ' '.join(slide_text)
        })

    return slides_content


def format_text_output(slides_content):
    """Format slides content as readable text."""
    output = []
    for slide in slides_content:
        output.append(f"\n{'='*60}")
        output.append(f"Slide {slide['slide_num']}")
        output.append('='*60)

        for text in slide['text']:
            output.append(text)

        if slide['tables']:
            output.append("\n[Tables]")
            for i, table in enumerate(slide['tables']):
                output.append(f"  Table {i+1}:")
                for row in table:
                    output.append(f"    {' | '.join(row)}")

    return '\n'.join(output)


def main():
    parser = argparse.ArgumentParser(description='Extract content from PPTX files')
    parser.add_argument('file', help='Path to .pptx file')
    parser.add_argument('--output', choices=['json', 'text'], default='text',
                        help='Output format (default: text)')
    parser.add_argument('--slide', type=int, help='Extract specific slide number')
    parser.add_argument('--encoding', default='utf-8', help='Output encoding')

    args = parser.parse_args()

    if not Path(args.file).exists():
        print(f"Error: File not found: {args.file}")
        sys.exit(1)

    try:
        content = extract_pptx(args.file, args.slide)

        if args.output == 'json':
            print(json.dumps(content, ensure_ascii=False, indent=2))
        else:
            print(format_text_output(content))

    except Exception as e:
        print(f"Error extracting content: {e}")
        sys.exit(1)


if __name__ == '__main__':
    main()
